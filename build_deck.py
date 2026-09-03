"""
EduNeuro — Pitch Deck PPTX Generator
Matches the pitch-deck.html design: near-black base, electric blue accent, 13 slides.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── palette (matching pitch-deck.html) ──────────────────────────────────
BG        = RGBColor(0x09, 0x09, 0x0B)   # near-black
BG2       = RGBColor(0x13, 0x13, 0x16)   # card bg
BG3       = RGBColor(0x1A, 0x1A, 0x1F)   # hover/input bg
WHITE     = RGBColor(0xFA, 0xFA, 0xFA)
TEXT2     = RGBColor(0xA1, 0xA1, 0xAA)
TEXT3     = RGBColor(0x71, 0x71, 0x7A)
ACCENT    = RGBColor(0x3B, 0x82, 0xF6)   # electric blue
ACCENT2   = RGBColor(0x60, 0xA5, 0xFA)
GREEN     = RGBColor(0x22, 0xC5, 0x5E)
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)
RED       = RGBColor(0xEF, 0x44, 0x44)
PURPLE    = RGBColor(0xA8, 0x55, 0xF7)
TEAL      = RGBColor(0x14, 0xB8, 0xA6)
BORDER    = RGBColor(0x12, 0x12, 0x14)   # subtle border

# ── helpers ──────────────────────────────────────────────────────────────

def set_run(run, text, size_pt, bold=False, color=WHITE, italic=False):
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text_box(slide, left, top, width, height, text_lines,
                 align=PP_ALIGN.LEFT, default_size=18, default_color=WHITE, default_bold=False):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if isinstance(line, str):
            run = p.add_run()
            set_run(run, line, default_size, bold=default_bold, color=default_color)
        else:
            run = p.add_run()
            sz = line.get("size", default_size)
            clr = line.get("color", default_color)
            bld = line.get("bold", default_bold)
            itl = line.get("italic", False)
            aln = line.get("align", None)
            if aln is not None:
                p.alignment = aln
            set_run(run, line["text"], sz, bold=bld, color=clr, italic=itl)
        p.space_after = Pt(4)
    return txb


def add_rect(slide, left, top, width, height, fill_rgb=None, line_rgb=None, line_pt=1):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def solid_bg(slide, rgb):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_rgb=rgb)


def slide_number(slide, n, total=15):
    add_text_box(slide, 12.2, 7.1, 1, 0.3,
                 [{"text": f"{n}/{total}", "size": 9, "color": TEXT3}],
                 align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════
#  BUILD
# ══════════════════════════════════════════════════════════════════════

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 1 — COVER
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

# logo area
add_text_box(s, 9.8, 0.35, 3.2, 0.55,
             [{"text": "EDUNEUR0", "size": 14, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.RIGHT)

# headline
add_text_box(s, 0.6, 1.4, 12, 2.0,
             [
                 {"text": "THE AI STUDY", "size": 54, "color": WHITE, "bold": True},
                 {"text": "COMPANION", "size": 54, "color": WHITE, "bold": True},
             ],
             align=PP_ALIGN.LEFT)

# accent divider
add_rect(s, 0.6, 3.35, 3.5, 0.06, fill_rgb=ACCENT)

# sub
add_text_box(s, 0.6, 3.55, 11, 0.9,
             [{"text": "AI-powered study orchestration for students",
               "size": 18, "color": TEXT2}],
             align=PP_ALIGN.LEFT)

# meta
add_text_box(s, 0.6, 4.8, 11, 0.8,
             [
                 {"text": "eduneuro.co.in   ·   September 2025   ·   Confidential",
                   "size": 11, "color": TEXT3},
             ],
             align=PP_ALIGN.LEFT)

slide_number(s, 1, 15)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 2 — MISSION
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

add_text_box(s, 0.6, 0.5, 4, 0.4,
             [{"text": "OUR MISSION", "size": 11, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.LEFT)
add_text_box(s, 0.6, 0.85, 12, 0.8,
             [{"text": "India's students deserve better than fragmented preparation tools.",
               "size": 28, "color": WHITE, "bold": True}],
             align=PP_ALIGN.LEFT)
slide_number(s, 2, 15)

# Three mission pillars
pillars = [
    ("🧠", "Student-First AI",
     "We don't just deliver content — we understand how each student learns, struggles, and improves, then adapt in real time."),
    ("📚", "Unified Learning Layer",
     "One platform that connects study planning, content, practice, doubt resolution, focus, and analytics — instead of 6 disconnected apps."),
    ("🇮🇳", "Built for India",
     "Designed for the realities of Indian competitive exams — GATE, JEE, NEET, UPSC — with curricula, languages, and patterns our students actually face."),
]

for i, (icon, title, desc) in enumerate(pillars):
    x = 0.6 + i * 4.2
    add_rect(s, x, 1.85, 4.0, 3.6, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
    add_rect(s, x, 1.85, 4.0, 0.06, fill_rgb=ACCENT)
    add_text_box(s, x + 0.2, 2.0, 3.6, 0.6,
                 [{"text": f"{icon}  {title}", "size": 15, "color": WHITE, "bold": True}])
    add_text_box(s, x + 0.2, 2.65, 3.6, 2.5,
                 [{"text": desc, "size": 12, "color": TEXT2}])

# Divider
add_rect(s, 0.6, 5.6, 12.1, 0.06, fill_rgb=ACCENT)
add_text_box(s, 0.6, 5.75, 12, 0.8,
             [{"text": "EduNeuro exists because the current study stack is broken — and we're fixing it, one student at a time.",
               "size": 14, "color": TEXT2, "italic": True}],
             align=PP_ALIGN.LEFT)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 3 — THE STORY / PROBLEM
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

add_text_box(s, 0.6, 0.5, 4, 0.4,
             [{"text": "THE STORY", "size": 11, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.LEFT)
slide_number(s, 3, 15)

# Persona box
add_rect(s, 0.6, 1.0, 12.1, 1.7, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
add_rect(s, 0.6, 1.0, 0.08, 1.7, fill_rgb=ACCENT)
add_text_box(s, 0.9, 1.1, 11.5, 1.5,
             [
                 {"text": "Meet Rahul — B.Tech CSE, GATE 2026 aspirant.",
                   "size": 15, "color": WHITE, "bold": True},
                 {"text": "Every day he opens 6 apps: YouTube, PDFs, question bank, Pomodoro timer, spreadsheet, ChatGPT. "
                          "By evening he's studied 8 hours but feels like he's going in circles.",
                   "size": 13, "color": TEXT2},
             ],
             align=PP_ALIGN.LEFT)

# Main headline
add_text_box(s, 0.6, 2.9, 12, 1.2,
             [
                 {"text": "Students don't need more resources.",
                   "size": 22, "color": WHITE, "bold": True},
                 {"text": "They need a system.",
                   "size": 22, "color": ACCENT, "bold": True},
             ],
             align=PP_ALIGN.LEFT)

# Tool chips (text-based)
tools = "YouTube   ·   Coaching   ·   PDFs   ·   ChatGPT   ·   Question Banks   ·   ✦  The Student  ✦   ·   Mock Tests   ·   Notes   ·   Timer   ·   Spreadsheet   ·   Peers"
add_text_box(s, 0.6, 4.2, 12, 0.6,
             [{"text": tools, "size": 11, "color": TEXT2}],
             align=PP_ALIGN.CENTER)

# Stats
for i, (num, lbl) in enumerate([("6–10", "Different tools per session"), ("0", "Systems that connect them")]):
    left = 2.5 + i * 5.5
    add_rect(s, left, 5.0, 4.5, 1.6, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
    add_text_box(s, left + 0.2, 5.1, 4.1, 0.7,
                 [{"text": num, "size": 36, "color": ACCENT, "bold": True}],
                 align=PP_ALIGN.LEFT)
    add_text_box(s, left + 0.2, 5.75, 4.1, 0.6,
                 [{"text": lbl, "size": 12, "color": TEXT3}],
                 align=PP_ALIGN.LEFT)

# Quote
add_rect(s, 0.6, 6.85, 12.1, 0.55, fill_rgb=BG2, line_rgb=ACCENT, line_pt=1)
add_text_box(s, 0.8, 6.9, 11.7, 0.5,
             [{"text": "Students are forced to become their own learning operating system.",
               "size": 12, "color": TEXT2, "italic": True}],
             align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 3 — THE INSIGHT
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

add_text_box(s, 0.6, 0.5, 4, 0.4,
             [{"text": "THE INSIGHT", "size": 11, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.LEFT)
add_text_box(s, 0.6, 0.85, 12, 1.0,
             [{"text": "The missing layer is student intelligence.",
               "size": 28, "color": WHITE, "bold": True}],
             align=PP_ALIGN.LEFT)
slide_number(s, 4, 15)
add_rect(s, 0.6, 2.0, 5.7, 4.5, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
add_text_box(s, 0.8, 2.1, 5.3, 0.5,
             [{"text": "✕  Today", "size": 14, "color": TEXT3, "bold": True}])
for j, item in enumerate(["Content platforms", "Tutoring services", "Planning tools",
                          "Practice platforms", "Focus apps", "Analytics dashboards"]):
    add_text_box(s, 0.9, 2.65 + j*0.5, 5.0, 0.45,
                 [{"text": item, "size": 12, "color": TEXT2}],
                 align=PP_ALIGN.LEFT)
add_text_box(s, 0.8, 5.6, 5.3, 0.5,
             [{"text": "All disconnected. Each knows one part.",
               "size": 11, "color": TEXT3, "italic": True}])

add_rect(s, 6.9, 2.0, 5.7, 4.5, fill_rgb=BG2, line_rgb=ACCENT, line_pt=2)
add_text_box(s, 7.1, 2.1, 5.3, 0.5,
             [{"text": "✓  EduNeuro", "size": 14, "color": ACCENT, "bold": True}])
for j, item in enumerate(["Student Intelligence Layer", "Full Study Loop Orchestration",
                          "Context-Aware Adaptation", "Performance-Driven Planning", "Personalized Next Action"]):
    add_text_box(s, 7.2, 2.65 + j*0.5, 5.0, 0.45,
                 [{"text": item, "size": 12, "color": WHITE}],
                 align=PP_ALIGN.LEFT)
add_text_box(s, 7.1, 5.5, 5.3, 0.5,
             [{"text": "One system. One student. Complete loop.",
               "size": 11, "color": ACCENT, "bold": True}])

add_text_box(s, 0.6, 6.7, 12, 0.8,
             [{"text": "Every tool knows one part of the workflow. EduNeuro is designed to understand the student.",
               "size": 14, "color": TEXT2}],
             align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 4 — WHY NOW
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

add_text_box(s, 0.6, 0.5, 4, 0.4,
             [{"text": "WHY NOW", "size": 11, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.LEFT)
add_text_box(s, 0.6, 0.85, 12, 1.0,
             [{"text": "The category is validating.", "size": 28, "color": WHITE, "bold": True}],
             align=PP_ALIGN.LEFT)
slide_number(s, 4, 15)

cards = [
    ("Duolingo Max", "AI-First Learning", "$1B", "Revenue · 130M MAU · 10.9M paying subs"),
    ("Khanmigo", "AI Tutor · Khan Academy", "65K+", "Students across 53 school districts"),
    ("Quizlet Q-Chat", "AI Study Companion", "Acquired", "Coconote (AI note-taking) · Feb 2026"),
    ("AI in Education", "Global Market", "$5.82B", "Projected by 2030 · 17.5% CAGR"),
]

for i, (name, typ, stat, label) in enumerate(cards):
    left = 0.6 + (i % 2) * 6.35
    top  = 2.0 + (i // 2) * 2.4
    add_rect(s, left, top, 6.0, 2.15, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
    add_rect(s, left, top, 6.0, 0.06, fill_rgb=ACCENT if i % 2 == 0 else TEAL)
    add_text_box(s, left + 0.2, top + 0.15, 5.5, 0.45,
                 [{"text": name, "size": 14, "color": WHITE, "bold": True}])
    add_text_box(s, left + 0.2, top + 0.55, 5.5, 0.4,
                 [{"text": typ, "size": 10, "color": TEXT3}])
    add_text_box(s, left + 0.2, top + 0.9, 5.5, 0.55,
                 [{"text": stat, "size": 20, "color": ACCENT, "bold": True}])
    add_text_box(s, left + 0.2, top + 1.4, 5.5, 0.55,
                 [{"text": label, "size": 11, "color": TEXT3}])

add_rect(s, 0.6, 6.8, 12.1, 0.55, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
add_text_box(s, 0.8, 6.85, 11.7, 0.5,
             [{"text": "Students are already paying for AI-powered learning companions. "
                       "Duolingo proved the subscription model works at scale.",
               "size": 12, "color": TEXT2}],
             align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 5 — PRODUCT ARCHITECTURE
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

add_text_box(s, 0.6, 0.5, 4, 0.4,
             [{"text": "PRODUCT ARCHITECTURE", "size": 11, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.LEFT)
add_text_box(s, 0.6, 0.85, 12, 1.0,
             [{"text": "One system for the complete study loop.", "size": 28, "color": WHITE, "bold": True}],
             align=PP_ALIGN.LEFT)
slide_number(s, 5, 15)

# Loop row
loop_items = ["PLAN", "LEARN", "PRACTICE", "FOCUS", "MEASURE", "ADAPT"]
total_w = len(loop_items)
item_w = 1.3
gap = 0.35
start_x = (13.33 - (total_w * item_w + (total_w - 1) * gap)) / 2

for i, label in enumerate(loop_items):
    x = start_x + i * (item_w + gap)
    add_rect(s, x, 2.8, item_w, 1.3, fill_rgb=BG2, line_rgb=ACCENT, line_pt=1)
    add_text_box(s, x, 3.05, item_w, 0.8,
                 [{"text": label, "size": 11, "color": WHITE, "bold": True}],
                 align=PP_ALIGN.CENTER)
    if i < len(loop_items) - 1:
        add_text_box(s, x + item_w + 0.05, 3.15, 0.25, 0.6,
                     [{"text": "→", "size": 14, "color": TEXT3}])

# Core circle (centered below)
add_rect(s, 5.4, 4.5, 2.5, 2.5, fill_rgb=BG, line_rgb=ACCENT, line_pt=3)
add_text_box(s, 5.4, 5.0, 2.5, 1.0,
             [{"text": "Student\nIntelligence", "size": 12, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.CENTER)

# Description
add_text_box(s, 0.6, 7.1, 12, 0.5,
             [{"text": "The platform continuously uses context, performance, behaviour and study history to personalize what comes next.",
               "size": 12, "color": TEXT2}],
             align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 6 — LIVE PRODUCT
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

add_text_box(s, 0.6, 0.5, 4, 0.4,
             [{"text": "LIVE PRODUCT", "size": 11, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.LEFT)
add_text_box(s, 0.6, 0.85, 12, 0.8,
             [{"text": "Built around the way students actually study.", "size": 22, "color": WHITE, "bold": True}],
             align=PP_ALIGN.LEFT)
slide_number(s, 6, 15)

features = [
    ("AI Study Planner", "Intelligent scheduling that adapts to your pace, goals, and performance patterns.", "Live"),
    ("AI Doubt Engine", "Context-aware tutoring that understands your syllabus, branch, and current topic.", "Live"),
    ("Virtual Library", "Collaborative study environment with rooms, focus timers, and peer presence.", "Live"),
    ("GATE Practice Engine", "Branch-aware practice with ECE/CSE/IN syllabi and year-wise trend analysis.", "Live"),
    ("Learning Analytics", "Real-time insights into study behaviour, topic mastery, and progress tracking.", "Live"),
    ("Focus & Pomodoro", "AI-powered distraction management and adaptive break scheduling.", "Building"),
]

for i, (title, desc, status) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.6 + col * 4.2
    y = 1.85 + row * 2.5
    add_rect(s, x, y, 3.9, 2.25, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
    add_rect(s, x, y, 3.9, 0.06, fill_rgb=ACCENT)
    add_text_box(s, x + 0.15, y + 0.12, 3.6, 0.5,
                 [{"text": title, "size": 13, "color": WHITE, "bold": True}])
    badge_color = ACCENT if status == "Live" else AMBER
    badge_text = status.upper()
    add_rect(s, x + 2.8, y + 0.13, 0.9, 0.28, fill_rgb=badge_color)
    add_text_box(s, x + 2.8, y + 0.13, 0.9, 0.28,
                 [{"text": badge_text, "size": 7, "color": BG, "bold": True}],
                 align=PP_ALIGN.CENTER)
    add_text_box(s, x + 0.15, y + 0.6, 3.6, 1.4,
                 [{"text": desc, "size": 11, "color": TEXT3}])

# ─────────────────────────────────────────────────────────────────────
# SLIDE 7 — TRACTION
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

add_text_box(s, 0.6, 0.5, 4, 0.4,
             [{"text": "TRACTION", "size": 11, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.LEFT)
add_text_box(s, 0.6, 0.85, 12, 0.8,
             [{"text": "Students are already showing up.", "size": 28, "color": WHITE, "bold": True}],
             align=PP_ALIGN.LEFT)
slide_number(s, 7, 15)

tractions = [
    ("505", "Visitors", "Unique visitors"),
    ("3,093", "Page Views", "Total engagement"),
    ("~400", "Peak Daily", "Daily active users"),
    ("250K+", "Organic Reach", "Social impressions"),
]

for i, (num, unit, lbl) in enumerate(tractions):
    col = i % 4
    x = 0.6 + col * 3.2
    add_rect(s, x, 1.9, 2.9, 3.0, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
    add_rect(s, x, 1.9, 2.9, 0.06, fill_rgb=ACCENT)
    add_text_box(s, x + 0.15, 2.1, 2.6, 1.0,
                 [{"text": num, "size": 40, "color": WHITE, "bold": True}],
                 align=PP_ALIGN.CENTER)
    add_text_box(s, x + 0.15, 3.05, 2.6, 0.4,
                 [{"text": unit, "size": 11, "color": ACCENT, "bold": True}],
                 align=PP_ALIGN.CENTER)
    add_text_box(s, x + 0.15, 3.45, 2.6, 0.6,
                 [{"text": lbl, "size": 10, "color": TEXT3}],
                 align=PP_ALIGN.CENTER)

# Context bar
add_rect(s, 0.6, 5.1, 12.1, 0.7, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
add_text_box(s, 0.8, 5.2, 11.7, 0.5,
             [{"text": "Open beta   ·   Organic acquisition   ·   Pre-revenue   ·   Strong early demand signals",
               "size": 12, "color": TEXT2}],
             align=PP_ALIGN.CENTER)

# Social proof
add_text_box(s, 0.6, 6.0, 12, 0.6,
             [{"text": "224,764 views   ·   +1,840 followers   ·   23,415 interactions",
               "size": 14, "color": WHITE, "bold": True}],
             align=PP_ALIGN.CENTER)
add_text_box(s, 0.6, 6.55, 12, 0.5,
             [{"text": "Social media analytics · eduneuro.co.in",
               "size": 10, "color": TEXT3}],
             align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 8 — USER SIGNAL
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

add_text_box(s, 0.6, 0.5, 4, 0.4,
             [{"text": "USER SIGNAL", "size": 11, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.LEFT)
add_text_box(s, 0.6, 0.85, 12, 0.8,
             [{"text": "Early users are not just curious.", "size": 28, "color": WHITE, "bold": True}],
             align=PP_ALIGN.LEFT)
slide_number(s, 8, 15)

# Big quote
add_rect(s, 0.6, 1.8, 12.1, 1.3, fill_rgb=BG2, line_rgb=ACCENT, line_pt=1)
add_text_box(s, 0.8, 1.95, 11.7, 1.0,
             [{"text": "100+ respondents indicated willingness to pay for deeper access to the platform — "
                       "before any monetization was introduced.",
               "size": 16, "color": TEXT2, "italic": True}],
             align=PP_ALIGN.CENTER)

# Two cards
add_rect(s, 0.6, 3.3, 5.7, 3.5, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
add_rect(s, 0.6, 3.3, 5.7, 0.06, fill_rgb=ACCENT)
add_text_box(s, 0.8, 3.45, 5.3, 0.5,
             [{"text": "Willingness to Pay", "size": 15, "color": WHITE, "bold": True}])
add_text_box(s, 0.8, 4.05, 5.3, 0.8,
             [{"text": "100+", "size": 42, "color": ACCENT, "bold": True}])
add_text_box(s, 0.8, 4.85, 5.3, 1.7,
             [{"text": "Survey respondents across BTech, MCA, ECE, CSE, Civil, Biotechnology and other branches "
                       "indicated monthly willingness to pay ranging from ₹100 to ₹500+ for platform access.",
               "size": 12, "color": TEXT3}])

add_rect(s, 7.0, 3.3, 5.7, 3.5, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
add_rect(s, 7.0, 3.3, 5.7, 0.06, fill_rgb=TEAL)
add_text_box(s, 7.2, 3.45, 5.3, 0.5,
             [{"text": "User Diversity", "size": 15, "color": WHITE, "bold": True}])
add_text_box(s, 7.2, 4.05, 5.3, 0.8,
             [{"text": "Multiple\nBranches", "size": 32, "color": ACCENT, "bold": True}])
add_text_box(s, 7.2, 4.85, 5.3, 1.7,
             [{"text": "Respondents span Electrical, CSE, IT, Civil, Biotechnology, ECE, and more — "
                       "indicating cross-branch appeal beyond a single exam category before any targeted marketing.",
               "size": 12, "color": TEXT3}])

# WTP breakdown
add_rect(s, 0.6, 7.0, 12.1, 0.4, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
add_text_box(s, 0.8, 7.05, 11.7, 0.35,
             [{"text": "WTP Distribution: ₹100–₹500/month   ·   Median: ₹325–₹350   ·   100+ of 126 willing to pay ≥ ₹100",
               "size": 10, "color": TEXT3}],
             align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 9 — MARKET
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

add_text_box(s, 0.6, 0.5, 4, 0.4,
             [{"text": "MARKET OPPORTUNITY", "size": 11, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.LEFT)
add_text_box(s, 0.6, 0.85, 12, 0.8,
             [{"text": "A focused wedge into a massive market.", "size": 28, "color": WHITE, "bold": True}],
             align=PP_ALIGN.LEFT)
slide_number(s, 9, 15)

layers = [
    ("GATE Preparation", "~1L registrations", "High-intent students actively investing in preparation", ACCENT),
    ("Competitive Exams", "50L+ candidates", "JEE, NEET, UPSC, CAT, SSC, Banking", TEAL),
    ("Student Learning", "4Cr+ students", "All Indian higher-education students", GREEN),
]

for i, (title, size, desc, color) in enumerate(layers):
    y = 1.8 + i * 1.5
    add_rect(s, 0.6, y, 12.1, 1.3, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
    add_rect(s, 0.6, y, 0.12, 1.3, fill_rgb=color)
    add_text_box(s, 0.9, y + 0.1, 4.0, 0.5,
                 [{"text": title, "size": 14, "color": WHITE, "bold": True}])
    add_text_box(s, 0.9, y + 0.55, 4.0, 0.6,
                 [{"text": desc, "size": 11, "color": TEXT3}])
    add_text_box(s, 9.5, y + 0.25, 3.0, 0.7,
                 [{"text": size, "size": 18, "color": ACCENT, "bold": True}],
                 align=PP_ALIGN.RIGHT)

# Bottom-up calc
add_text_box(s, 0.6, 6.3, 12, 0.4,
             [{"text": "Addressable Market — Bottom-Up Calculation (GATE)", "size": 11, "color": TEXT3, "bold": True}],
             align=PP_ALIGN.LEFT)

calc_rows = [
    ("GATE serious candidates", "~50,000 students", TEXT2),
    ("× Avg WTP / year", "₹3,600 (₹300/mo avg)", TEXT2),
    ("GATE-only addressable", "~₹18 Cr", ACCENT),
    ("Cross-exam expansion", "₹100s of Crores", ACCENT),
]

for i, (lbl, val, clr) in enumerate(calc_rows):
    y = 6.75 + i * 0.38
    is_total = i >= 2
    if is_total:
        add_rect(s, 0.6, y, 12.1, 0.38, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
    add_text_box(s, 0.8, y + 0.02, 8, 0.35,
                 [{"text": lbl, "size": 10, "color": WHITE if is_total else TEXT3}])
    add_text_box(s, 10.0, y + 0.02, 2.5, 0.35,
                 [{"text": val, "size": 10, "color": clr, "bold": True}],
                 align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 10 — COMPETITION
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

add_text_box(s, 0.6, 0.5, 4, 0.4,
             [{"text": "COMPETITIVE LANDSCAPE", "size": 11, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.LEFT)
add_text_box(s, 0.6, 0.85, 12, 0.8,
             [{"text": "Not competing with everyone. Connecting the whole loop.", "size": 24, "color": WHITE, "bold": True}],
             align=PP_ALIGN.LEFT)
slide_number(s, 10, 15)

add_text_box(s, 0.6, 1.75, 12, 0.4,
             [{"text": "No one is optimising the full study loop. That is the opening.",
               "size": 13, "color": TEXT2}],
             align=PP_ALIGN.LEFT)

competitors = [
    ("Unacademy / PW / Testbook", "Content + Coaching", "Content libraries, live classes, mock tests. Strong on delivery, weak on personalisation.", False),
    ("ChatGPT / Gemini", "Open-ended AI Tutoring", "Reasoning, general knowledge, open Q&A. Not exam-specialised, no curriculum map.", False),
    ("Notion / Forest", "Study Productivity", "Notes & focus timers. No adaptive engine, no exam-specific intelligence.", False),
    ("EduNeuro", "Study Orchestration", "Full-loop adaptive study platform: plans → content → practice → focus → analytics — personalised per student.", True),
]

for i, (name, cat, desc, is_us) in enumerate(competitors):
    y = 2.3 + i * 1.15
    bg = BG if is_us else BG2
    border = ACCENT if is_us else BORDER
    add_rect(s, 0.6, y, 12.1, 1.05, fill_rgb=bg, line_rgb=border, line_pt=(2 if is_us else 1))
    add_text_box(s, 0.8, y + 0.12, 3.0, 0.8,
                 [{"text": name, "size": 13, "color": ACCENT if is_us else WHITE, "bold": True}])
    add_text_box(s, 4.0, y + 0.12, 2.5, 0.8,
                 [{"text": cat, "size": 11, "color": TEXT3}])
    add_text_box(s, 6.7, y + 0.12, 5.8, 0.8,
                 [{"text": desc, "size": 11, "color": TEXT2}])

add_rect(s, 0.6, 6.95, 12.1, 0.5, fill_rgb=BG2, line_rgb=ACCENT, line_pt=1)
add_text_box(s, 0.8, 7.0, 11.7, 0.4,
             [{"text": "EduNeuro is not trying to beat every competitor. It is trying to connect the entire study workflow around the individual student.",
               "size": 11, "color": TEXT2}],
             align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 11 — GTM
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

add_text_box(s, 0.6, 0.5, 4, 0.4,
             [{"text": "GO-TO-MARKET", "size": 11, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.LEFT)
add_text_box(s, 0.6, 0.85, 12, 0.8,
             [{"text": "Start focused. Expand systematically.", "size": 28, "color": WHITE, "bold": True}],
             align=PP_ALIGN.LEFT)
slide_number(s, 11, 15)

gtm = [
    ("Phase 1 · Now", "GATE Preparation",
     "Deep integration with GATE syllabi, branch-specific practice, trend analysis, and study planning for CSE, ECE, IN. High-intent users who actively invest in preparation.",
     "1L+", "Registrations", "~50K", "Serious"),
    ("Phase 2 · 6–12 months", "Expansion Exams",
     "Leverage the same student intelligence layer for JEE, NEET, UPSC, CAT. The orchestration engine is exam-agnostic.",
     "50L+", "Combined", "Multi-exam", "Same Core"),
]

for idx, (pill, title, desc, s1v, s1l, s2v, s2l) in enumerate(gtm):
    x = 0.6 + idx * 6.35
    add_rect(s, x, 1.85, 6.0, 4.0, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
    add_rect(s, x, 1.85, 6.0, 0.06, fill_rgb=ACCENT)
    add_rect(s, x + 0.2, 2.0, 2.2, 0.4, fill_rgb=ACCENT)
    add_text_box(s, x + 0.2, 2.0, 2.2, 0.4,
                 [{"text": pill, "size": 9, "color": BG, "bold": True}],
                 align=PP_ALIGN.CENTER)
    add_text_box(s, x + 0.2, 2.55, 5.5, 0.6,
                 [{"text": title, "size": 18, "color": WHITE, "bold": True}])
    add_text_box(s, x + 0.2, 3.2, 5.5, 2.0,
                 [{"text": desc, "size": 12, "color": TEXT2}])
    for j, (v, l) in enumerate([(s1v, s1l), (s2v, s2l)]):
        add_rect(s, x + 0.2 + j * 2.6, 5.3, 2.4, 1.2, fill_rgb=BG3)
        add_text_box(s, x + 0.3 + j * 2.6, 5.35, 2.2, 0.6,
                     [{"text": v, "size": 16, "color": ACCENT, "bold": True}])
        add_text_box(s, x + 0.3 + j * 2.6, 5.9, 2.2, 0.5,
                     [{"text": l, "size": 10, "color": TEXT3}],
                     align=PP_ALIGN.CENTER)

# Channels
add_rect(s, 0.6, 6.1, 12.1, 0.5, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
add_text_box(s, 0.8, 6.15, 11.7, 0.4,
             [{"text": "Organic Social   ·   Content Marketing   ·   Community   ·   University Outreach   ·   Student Ambassadors   ·   Referral Program",
               "size": 11, "color": TEXT2}],
             align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 12 — ROADMAP
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

add_text_box(s, 0.6, 0.5, 4, 0.4,
             [{"text": "ROADMAP", "size": 11, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.LEFT)
add_text_box(s, 0.6, 0.85, 12, 0.8,
             [{"text": "Building the student intelligence layer.", "size": 28, "color": WHITE, "bold": True}],
             align=PP_ALIGN.LEFT)
slide_number(s, 12, 15)

phases = [
    ("NOW", "2025", [
        "GATE CSE / ECE / IN",
        "AI Study Planner",
        "Virtual Library",
        "AI Doubt Engine",
        "Organic Growth",
    ], True),
    ("6 MO", "2026", [
        "JEE Integration",
        "NEET Module",
        "Focus & Pomodoro",
        "Monetization",
    ], False),
    ("12 MO", "2026", [
        "UPSC / CAT",
        "Advanced Analytics",
        "Peer Learning",
        "Mobile App",
    ], False),
    ("18+ MO", "2027", [
        "AI Tutoring",
        "Adaptive Engine",
        "Multi-language",
        "Platform API",
    ], False),
]

for i, (label, period, items, is_now) in enumerate(phases):
    x = 0.6 + i * 3.2
    border_c = ACCENT if is_now else BORDER
    fill_top = BG2
    add_rect(s, x, 1.85, 3.0, 4.2, fill_rgb=fill_top, line_rgb=border_c, line_pt=(2 if is_now else 1))
    add_rect(s, x, 1.85, 3.0, 0.06, fill_rgb=ACCENT if is_now else BORDER)
    add_text_box(s, x, 2.0, 3.0, 0.5,
                 [{"text": label, "size": 14, "color": ACCENT if is_now else WHITE, "bold": True}],
                 align=PP_ALIGN.CENTER)
    add_text_box(s, x, 2.45, 3.0, 0.4,
                 [{"text": period, "size": 11, "color": TEXT3}],
                 align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        item_color = ACCENT if is_now else TEXT2
        add_rect(s, x + 0.15, 3.1 + j * 0.75, 2.7, 0.6, fill_rgb=BG3, line_rgb=border_c, line_pt=1)
        add_text_box(s, x + 0.25, 3.15 + j * 0.75, 2.5, 0.5,
                     [{"text": item, "size": 10, "color": item_color}])

# ─────────────────────────────────────────────────────────────────────
# SLIDE 14 — TEAM
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

add_text_box(s, 0.6, 0.5, 4, 0.4,
             [{"text": "TEAM", "size": 11, "color": ACCENT, "bold": True}],
             align=PP_ALIGN.LEFT)
add_text_box(s, 0.6, 0.85, 12, 0.8,
             [{"text": "The right team for the right problem.", "size": 28, "color": WHITE, "bold": True}],
             align=PP_ALIGN.LEFT)
slide_number(s, 14, 15)

# Founder highlight
add_rect(s, 0.6, 1.85, 12.1, 1.3, fill_rgb=BG2, line_rgb=ACCENT, line_pt=2)
add_text_box(s, 0.9, 2.0, 11.5, 0.9,
             [{"text": "Sumanta Bhargab  ·  Founder & CEO",
               "size": 16, "color": WHITE, "bold": True}])
add_text_box(s, 0.9, 2.7, 11.5, 0.4,
             [{"text": "Building EduNeuro from first-principles understanding of how students actually learn. Self-taught AI engineer, GATE-qualified, driving a product that serves the next generation of Indian students.",
               "size": 12, "color": TEXT2}])

# Advisors / mentors placeholder
add_text_box(s, 0.6, 3.35, 12, 0.5,
             [{"text": "Building the core team. Open to exceptional advisors, educators, and technologists who share the vision.",
               "size": 12, "color": TEXT3}],
             align=PP_ALIGN.LEFT)

# Stats bar
add_rect(s, 0.6, 4.1, 12.1, 1.0, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
for i, (num, lbl) in enumerate([
    ("GATE Qualified", "Technical depth"),
    ("AI-First", "Built from scratch"),
    ("Self-Funded", "Bootstrap conviction"),
]):
    x = 0.9 + i * 4.0
    add_text_box(s, x, 4.2, 3.5, 0.35,
                 [{"text": num, "size": 14, "color": ACCENT, "bold": True}])
    add_text_box(s, x, 4.55, 3.5, 0.35,
                 [{"text": lbl, "size": 10, "color": TEXT3}])

# Hiring / joining message
add_rect(s, 0.6, 5.35, 12.1, 0.55, fill_rgb=BG2, line_rgb=ACCENT, line_pt=1)
add_text_box(s, 0.8, 5.4, 11.7, 0.45,
             [{"text": "We are hiring exceptional engineers, designers, and educators who want to build the future of learning.",
               "size": 12, "color": TEXT2}],
             align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────────────────────────────────
# SLIDE 15 — CLOSING
# ─────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6])
solid_bg(s, BG)

# Big background number
add_text_box(s, 0.3, 0.5, 13, 3.5,
             [{"text": "4Cr+", "size": 200, "color": RGBColor(0x10, 0x14, 0x28), "bold": True}],
             align=PP_ALIGN.LEFT)

add_text_box(s, 0.6, 1.6, 12, 1.2,
             [{"text": "students. One platform. Infinite potential.",
               "size": 38, "color": WHITE, "bold": True}],
             align=PP_ALIGN.LEFT)

add_rect(s, 0.6, 2.9, 3.5, 0.06, fill_rgb=ACCENT)

add_text_box(s, 0.6, 3.1, 11, 0.9,
             [{"text": "EduNeuro is building the operating system for how India's students learn.",
               "size": 20, "color": TEXT2}],
             align=PP_ALIGN.LEFT)

# Three pillars
pillars = [
    ("🧠", "AI-NATIVE", "Adaptive by design — not bolted on"),
    ("📱", "MOBILE-FIRST", "Built for the smartphone generation"),
    ("🇮🇳", "INDIA-DEEP", "Curricula, languages, exam patterns — local at every layer"),
]
for i, (icon, title, desc) in enumerate(pillars):
    x = 0.6 + i * 4.0
    add_rect(s, x, 4.3, 3.6, 1.8, fill_rgb=BG2, line_rgb=BORDER, line_pt=1)
    add_rect(s, x, 4.3, 3.6, 0.06, fill_rgb=ACCENT)
    add_text_box(s, x + 0.2, 4.45, 3.2, 0.55,
                 [{"text": f"{icon}  {title}", "size": 16, "color": WHITE, "bold": True}])
    add_text_box(s, x + 0.2, 4.95, 3.2, 0.9,
                 [{"text": desc, "size": 13, "color": TEXT3}])

# Contact bar
add_rect(s, 0, 6.85, 13.33, 0.65, fill_rgb=BG2)
add_text_box(s, 0.5, 6.92, 12, 0.5,
             [{"text": "eduneuro.co.in  ·  hello@eduneuro.co.in  ·  linkedin.com/company/eduneuro  ·  @eduneuro_ai",
               "size": 11, "color": ACCENT}],
             align=PP_ALIGN.CENTER)

# Footer
add_text_box(s, 0.5, 7.3, 12, 0.2,
             [{"text": "Confidential · For discussion purposes only · EduNeuro Technologies Pvt. Ltd. · September 2025",
               "size": 8, "color": TEXT3}],
             align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════
#  SAVE
# ══════════════════════════════════════════════════════════════════════
import os
out_dir = "C:/Users/Sumanta Bhargab/eduniche/eduniche"
os.makedirs(out_dir, exist_ok=True)
out_path = os.path.join(out_dir, "EduNeuro_Pitch_Deck_v2.pptx")
prs.save(out_path)
print(f"[OK] Saved -> {out_path}")
print(f"    Slides: {len(prs.slides)}")
